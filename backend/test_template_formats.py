from docx import Document
from pptx import Presentation
from parser import TemplateParser
import os


def create_sample_docx(path: str):
    doc = Document()
    doc.add_paragraph("KPMG SAMPLE HEADER")
    doc.add_paragraph("Professional Summary")
    doc.save(path)


def create_sample_pptx(path: str):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "ZS Sample Slide"
    slide.placeholders[1].text = "Professional Experience"
    prs.save(path)


if __name__ == "__main__":
    base = os.path.dirname(__file__)
    docx_path = os.path.join(base, "_sample_template.docx")
    pptx_path = os.path.join(base, "_sample_template.pptx")

    create_sample_docx(docx_path)
    create_sample_pptx(pptx_path)

    docx_info = TemplateParser.extract_template_info(docx_path)
    pptx_info = TemplateParser.extract_template_info(pptx_path)

    print("DOCX sections:", len(docx_info.get("sections", [])))
    print("PPTX sections:", len(pptx_info.get("sections", [])))

    os.remove(docx_path)
    os.remove(pptx_path)
